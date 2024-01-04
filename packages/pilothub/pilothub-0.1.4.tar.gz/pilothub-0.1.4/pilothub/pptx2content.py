from pptx import Presentation


class PPTxFile:
    def __init__(self, file_path):
        self.prs = Presentation(file_path)
        self.slides = self.prs.slides
        self.notes = self.get_notes()
    
    def get_slide_text(self, slide):
        """Returns the text in a slide
        
        Arguments:
            slide {pptx.slide} -- slide object
        Returns:
            str -- text in slide
        """
        slide_text = '\n'.join([shape.text for shape in slide.shapes if 
                                hasattr(shape, "text")])
        return slide_text

    def get_slide_notes(self, slide):
        """Returns the notes in a slide
        
        Arguments:
            slide {pptx.slide} -- slide object
        Returns:
            str -- notes in slide
        """
        notes_slide = self.slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text
        return notes_text

    def set_slide_notes(self, slide, text):
        """Sets the notes in a slide
        
        Arguments:
            slide {pptx.slide} -- slide object
            text {str} -- text to be set as notes
        """
        self.slide.notes_slide.notes_text_frame.text = text

    def erase_slide_notes(self, slide):
        slide.notes_slide.notes_text_frame.clear()


















def process_presentation(file_path, output_path):
    prs = Presentation(file_path)
    ppt_len = len(prs.slides)

    for i,slide in enumerate(prs.slides):
        slide_text = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])

        ignore_types = ["CoverPage","Quote Slide","Agenda","Section Header","RunningMan-Infographic",
                        "QuoteHead"]
        
        notes_slide = slide.notes_slide
        # Clear existing text if any, and add the new text
        notes_slide.notes_text_frame.clear()

        if slide.slide_layout.name in ignore_types:
            notes_slide.notes_text_frame.text = ""            
        elif "quiz" in slide_text.lower():
            notes_slide.notes_text_frame.text = ""
        elif i>ppt_len-4:
            notes_slide.notes_text_frame.text = ""
        else:
            try:
                notes_text = get_notes(slide_text)
                notes_slide.notes_text_frame.text = notes_text
            except:
                notes_slide.notes_text_frame.text = slide_text

    prs.save(output_path)

# Example usage
file_path = r"C:\Users\Anshu Pandey\Downloads\PPT_001_Introduction to Big Data Ingestion.pptx"  # Replace with your input file path
output_path = 'PPT1_2.pptx'  # Replace with your output file path

# Uncomment the following line to run the function with your file paths
process_presentation(file_path, output_path)
